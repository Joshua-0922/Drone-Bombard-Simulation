#!/usr/bin/env python3
"""Drone-Bombard-Simulation Round 1 발표 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK     = RGBColor(0x1B, 0x1B, 0x2F)   # 슬라이드 배경
BG_CARD     = RGBColor(0x27, 0x27, 0x44)   # 카드 배경
ACCENT      = RGBColor(0x5C, 0xB8, 0xFF)   # 파란 강조
ACCENT2     = RGBColor(0xFF, 0x6B, 0x6B)   # 빨간 강조 (문제/실패)
ACCENT3     = RGBColor(0x51, 0xE0, 0x98)   # 초록 강조 (해결/성공)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_GRAY  = RGBColor(0xDD, 0xDD, 0xEE)
YELLOW      = RGBColor(0xFF, 0xD9, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # rounded corner
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6), font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, clr, bld) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = spacing
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=ACCENT):
    """data: list of lists. First row = header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "맑은 고딕"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                    paragraph.alignment = PP_ALIGN.CENTER

            # cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r == 0:
                cell_fill.fore_color.rgb = RGBColor(0x35, 0x35, 0x55)
            elif r % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x48)
            else:
                cell_fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3C)

            # border
            from pptx.oxml.ns import qn
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    from lxml import etree
                    ln = etree.SubElement(tcPr, qn(border_name))
                ln.set('w', '6350')
                sf = ln.find(qn('a:solidFill'))
                if sf is None:
                    from lxml import etree
                    sf = etree.SubElement(ln, qn('a:solidFill'))
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is None:
                    from lxml import etree
                    srgb = etree.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', '3A3A5C')

    return table_shape


# ════════════════════════════════════════════════════════════════
#  SLIDE 1: 표지
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
         "Drone-Bombard-Simulation", 44, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
         "강화학습 기반 드론 정밀투하 — 현황 보고", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "2026-05-25", 20, GRAY, False, PP_ALIGN.CENTER)

# 하단 장식선
add_rect(slide, Inches(4), Inches(5.0), Inches(5.3), Pt(3), ACCENT)


# ════════════════════════════════════════════════════════════════
#  SLIDE 2: 프로젝트 개요
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "프로젝트 개요", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(3), Pt(3), ACCENT)

# 목표 카드
add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.2), BG_CARD)
add_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.4),
         "목표", 20, ACCENT, True)
add_bullet_text(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(1.5), [
    ("ROS2 + Gazebo + PX4 시뮬레이션 환경", WHITE, False),
    ("SAC 강화학습으로 드론 정밀투하 정책 학습", WHITE, False),
    ("고정 타겟 (11m, 10m) ENU 좌표", GRAY, False),
], 15)

# 환경 카드
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.2), BG_CARD)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.4),
         "시뮬레이션 환경", 20, ACCENT, True)
add_bullet_text(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(1.5), [
    ("Gazebo Harmonic (Docker)", WHITE, False),
    ("PX4 SITL 비행제어", WHITE, False),
    ("Stable-Baselines3 SAC", WHITE, False),
], 15)

# 에이전트 카드
add_rect(slide, Inches(0.6), Inches(3.9), Inches(5.8), Inches(3.0), BG_CARD)
add_text(slide, Inches(0.9), Inches(4.0), Inches(5.2), Inches(0.4),
         "에이전트 설계", 20, ACCENT, True)
add_bullet_text(slide, Inches(0.9), Inches(4.5), Inches(5.2), Inches(2.2), [
    ("관측 (17차원)", YELLOW, True),
    ("  위치, 속도, 각속도, 타겟 상대좌표", GRAY, False),
    ("  CCIP 예측 낙하점, 탄착 예상시간", GRAY, False),
    ("행동 (5차원)", YELLOW, True),
    ("  vx, vy, vz, yaw rate, 투하 결정", GRAY, False),
], 14)

# 학습 규모 카드
add_rect(slide, Inches(6.8), Inches(3.9), Inches(5.8), Inches(3.0), BG_CARD)
add_text(slide, Inches(7.1), Inches(4.0), Inches(5.2), Inches(0.4),
         "학습 규모", 20, ACCENT, True)
add_bullet_text(slide, Inches(7.1), Inches(4.5), Inches(5.2), Inches(2.2), [
    ("150k~200k step / 약 1,000~1,500 에피소드", WHITE, False),
    ("에피소드 1회 = 최대 500 step (~25초)", WHITE, False),
    ("20Hz 제어 루프", GRAY, False),
], 15)


# ════════════════════════════════════════════════════════════════
#  SLIDE 3: 겪었던 문제들 — 개요
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "어떤 문제들을 겪었는가", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(5), Pt(3), ACCENT)

add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "총 7회 학습을 거치며 세 가지 유형의 실패를 경험했다.", 18, GRAY)

# 3개 카드
card_w = Inches(3.7)
card_h = Inches(4.8)
gap = Inches(0.5)
start_x = Inches(0.6)
card_y = Inches(2.1)

# 카드 1
add_rect(slide, start_x, card_y, card_w, card_h, BG_CARD, ACCENT2)
add_text(slide, start_x + Inches(0.3), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.5),
         "유형 1: 인프라 결함", 20, ACCENT2, True)
add_bullet_text(slide, start_x + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.6), Inches(3.5), [
    ("CCIP 예측 1m vs 실제 14.87m", WHITE, False),
    ("→ 14m 체계적 오차", ACCENT2, True),
    ("", WHITE, False),
    ("원인: spin thread crash로", GRAY, False),
    ("ROS2 콜백 freeze, CCIP 오염", GRAY, False),
    ("", WHITE, False),
    ("해결: jekyun_v2 인프라 패치", ACCENT3, False),
    ("→ 오차 해소 확인 ✓", ACCENT3, False),
], 14, spacing=Pt(4))

# 카드 2
x2 = start_x + card_w + gap
add_rect(slide, x2, card_y, card_w, card_h, BG_CARD, ACCENT2)
add_text(slide, x2 + Inches(0.3), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.5),
         "유형 2: Critic 발산", 20, ACCENT2, True)
add_bullet_text(slide, x2 + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.6), Inches(3.5), [
    ("critic loss 17,000 (정상 ~100)", WHITE, False),
    ("→ 정책 완전 붕괴", ACCENT2, True),
    ("", WHITE, False),
    ("원인: gradient_steps=4와", GRAY, False),
    ("큰 터미널 보상(+350) 충돌", GRAY, False),
    ("", WHITE, False),
    ("해결: gradient_steps 4 → 1", ACCENT3, False),
    ("→ critic 안정화 확인 ✓", ACCENT3, False),
], 14, spacing=Pt(4))

# 카드 3
x3 = x2 + card_w + gap
add_rect(slide, x3, card_y, card_w, card_h, BG_CARD, YELLOW)
add_text(slide, x3 + Inches(0.3), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.5),
         "유형 3: 정밀도 정체 ★", 20, YELLOW, True)
add_bullet_text(slide, x3 + Inches(0.3), card_y + Inches(0.9), card_w - Inches(0.6), Inches(3.5), [
    ("drop은 발생 (220회/112k step)", WHITE, False),
    ("→ 오차 13m에서 수렴, 개선 없음", ACCENT2, True),
    ("→ 성공(≤0.5m) 0회", ACCENT2, True),
    ("", WHITE, False),
    ("근본 원인 4가지:", YELLOW, True),
    ("① 보상이 정밀도에 둔감", WHITE, False),
    ("② 보상 스케일 과대", WHITE, False),
    ("③ auto-drop이 학습 차단", WHITE, False),
    ("④ 종료 조건 부족", WHITE, False),
], 14, spacing=Pt(4))


# ════════════════════════════════════════════════════════════════
#  SLIDE 4: 보상 함수 구조 (4-Layer 개요)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "보상 함수 — 4-Layer 구조", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(5), Pt(3), ACCENT)

# 4개 layer 카드 가로 배치
lw = Inches(2.85)
lh = Inches(5.0)
lx = Inches(0.6)
ly = Inches(1.4)
lgap = Inches(0.3)

# Layer 1
add_rect(slide, lx, ly, lw, lh, BG_CARD)
add_text(slide, lx + Inches(0.2), ly + Inches(0.15), lw - Inches(0.4), Inches(0.4),
         "Layer 1: 안전", 18, ACCENT2, True)
add_text(slide, lx + Inches(0.2), ly + Inches(0.55), lw - Inches(0.4), Inches(0.3),
         "매 step · 에피소드 종료", 12, GRAY)
add_bullet_text(slide, lx + Inches(0.2), ly + Inches(1.0), lw - Inches(0.4), Inches(3.5), [
    ("crash    -50", WHITE, False),
    ("overspeed  -30", WHITE, False),
    ("뒤집힘    -30", WHITE, False),
    ("", WHITE, False),
    ("스케일 축소 ↓", ACCENT3, True),
    ("이전 -100~-50", GRAY, False),
    ("현재  -50~-30", ACCENT3, False),
], 13, spacing=Pt(3))

# Layer 2
lx2 = lx + lw + lgap
add_rect(slide, lx2, ly, lw, lh, BG_CARD)
add_text(slide, lx2 + Inches(0.2), ly + Inches(0.15), lw - Inches(0.4), Inches(0.4),
         "Layer 2: 안정", 18, GRAY, True)
add_text(slide, lx2 + Inches(0.2), ly + Inches(0.55), lw - Inches(0.4), Inches(0.3),
         "매 step · 변경 없음", 12, GRAY)
add_bullet_text(slide, lx2 + Inches(0.2), ly + Inches(1.0), lw - Inches(0.4), Inches(3.5), [
    ("시간 지연    -0.05", WHITE, False),
    ("기체 불안정  -0.05", WHITE, False),
    ("급격 조작    -0.05", WHITE, False),
    ("", WHITE, False),
    ("변경 없음", GRAY, False),
    ("안정적으로 동작 중", GRAY, False),
], 13, spacing=Pt(3))

# Layer 3
lx3 = lx2 + lw + lgap
add_rect(slide, lx3, ly, lw, lh, BG_CARD)
add_text(slide, lx3 + Inches(0.2), ly + Inches(0.15), lw - Inches(0.4), Inches(0.4),
         "Layer 3: 접근", 18, GRAY, True)
add_text(slide, lx3 + Inches(0.2), ly + Inches(0.55), lw - Inches(0.4), Inches(0.3),
         "매 step · 변경 없음", 12, GRAY)
add_bullet_text(slide, lx3 + Inches(0.2), ly + Inches(1.0), lw - Inches(0.4), Inches(3.5), [
    ("거리 감소 보상  +0.5", WHITE, False),
    ("방향 정렬      +0.7", WHITE, False),
    ("CCIP 정확도    +0.4", WHITE, False),
    ("", WHITE, False),
    ("speed_gate로", GRAY, False),
    ("제자리 보상 축적 차단", GRAY, False),
    ("", WHITE, False),
    ("신호 자체는 정상.", GRAY, False),
    ("Layer 4가 덮어버린 게 문제", YELLOW, False),
], 13, spacing=Pt(3))

# Layer 4
lx4 = lx3 + lw + lgap
add_rect(slide, lx4, ly, lw, lh, BG_CARD, YELLOW)
add_text(slide, lx4 + Inches(0.2), ly + Inches(0.15), lw - Inches(0.4), Inches(0.4),
         "Layer 4: 투하 ★", 18, YELLOW, True)
add_text(slide, lx4 + Inches(0.2), ly + Inches(0.55), lw - Inches(0.4), Inches(0.3),
         "drop 시점 · 핵심 변경", 12, YELLOW)
add_bullet_text(slide, lx4 + Inches(0.2), ly + Inches(1.0), lw - Inches(0.4), Inches(3.5), [
    ("proximity bonus", ACCENT3, True),
    ("precision (k2 강화)", ACCENT3, True),
    ("예측 정확도 bonus", ACCENT3, True),
    ("jackpot", WHITE, False),
    ("", WHITE, False),
    ("이번 Round 1의", YELLOW, True),
    ("핵심 변경 대상", YELLOW, True),
    ("→ 다음 슬라이드 상세", YELLOW, False),
], 13, spacing=Pt(3))

add_text(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
         "Layer 1~3은 매 step, Layer 4는 투하 시점에 한 번 부여", 14, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  SLIDE 5: Layer 4 이전 설계와 문제
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "투하 보상 — 이전 설계의 문제", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(6), Pt(3), ACCENT)

# 수식
add_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(1.4), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.5),
         "이전 수식", 16, GRAY, True)
add_text(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(0.5),
         "R4  =  150  +  100 × exp(-0.3 × d_error)  +  100 × [d_error ≤ 0.1m]", 20, WHITE, True, font_name="Consolas")

# 문제 설명
add_rect(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.6), BG_CARD, ACCENT2)
add_text(slide, Inches(0.9), Inches(3.3), Inches(5.2), Inches(0.4),
         "무엇이 문제였나", 20, ACCENT2, True)
add_bullet_text(slide, Inches(0.9), Inches(3.9), Inches(5.2), Inches(2.8), [
    ("bonus 150이 정액 → drop만 하면 무조건 지급", WHITE, False),
    ("", WHITE, False),
    ("k2=0.3 → decay가 완만해서", WHITE, False),
    ("15m drop과 5m drop의 차이가 겨우 21", WHITE, False),
    ("", WHITE, False),
    ("이미 150을 확보한 에이전트 입장에서", GRAY, False),
    ("21을 더 얻으려 위험을 감수할 이유가 없다", YELLOW, True),
    ("", WHITE, False),
    ("터미널 범위 -100~+350 →", WHITE, False),
    ("매 step 접근 신호를 critic이 무시", ACCENT2, False),
], 14, spacing=Pt(3))

# 보상 테이블
data = [
    ["d_error", "보상", "비고"],
    ["0m", "350", "jackpot 포함"],
    ["5m", "172", ""],
    ["10m", "155", ""],
    ["15m", "151", "← 학습 수렴점"],
    ["50m", "150", "← bonus만으로 충분"],
]
add_table(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.0),
          6, 3, data,
          col_widths=[Inches(1.5), Inches(1.5), Inches(2.8)])

add_text(slide, Inches(6.8), Inches(6.4), Inches(5.8), Inches(0.5),
         "0m(350)과 50m(150)의 차이: 겨우 2.3배", 16, ACCENT2, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  SLIDE 6: Layer 4 현재 설계
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "투하 보상 — Round 1 재설계", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(5), Pt(3), ACCENT)

# 수식
add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.6), BG_CARD, ACCENT3)
add_text(slide, Inches(1.0), Inches(1.4), Inches(11), Inches(0.4),
         "현재 수식", 16, ACCENT3, True)
add_bullet_text(slide, Inches(1.0), Inches(1.85), Inches(11), Inches(1.0), [
    ("R4  =  30 × exp(-0.3 × d_xy)                           proximity bonus", WHITE, False),
    ("     + 100 × exp(-0.5 × d_error)                        정밀도", WHITE, False),
    ("     + 20 × exp(-0.1 × |d_impact − d_error|)     예측 정확도", WHITE, False),
    ("     + 50 × [d_error ≤ 0.1m]                              jackpot", WHITE, False),
], 14, font_name="Consolas", spacing=Pt(2))

# 3가지 변경
c_w = Inches(3.7)
c_h = Inches(3.4)
cx = Inches(0.6)
cy = Inches(3.3)
cg = Inches(0.5)

# 변경 1
add_rect(slide, cx, cy, c_w, c_h, BG_CARD)
add_text(slide, cx + Inches(0.2), cy + Inches(0.15), c_w - Inches(0.4), Inches(0.4),
         "① Bonus → 거리 의존형", 16, ACCENT3, True)
add_bullet_text(slide, cx + Inches(0.2), cy + Inches(0.65), c_w - Inches(0.4), Inches(2.5), [
    ("정액 150 → 30×exp(-0.3×d_xy)", WHITE, False),
    ("", WHITE, False),
    ("d_xy 50m → bonus ≈ 0", GRAY, False),
    ("d_xy 14m → bonus = 0.5", GRAY, False),
    ("d_xy  5m → bonus = 6.7", WHITE, False),
    ("d_xy  3m → bonus = 12.2", WHITE, False),
    ("d_xy  0m → bonus = 30", ACCENT3, True),
    ("", WHITE, False),
    ("가까이 가야 보상을 받는다", YELLOW, True),
], 13, spacing=Pt(2))

# 변경 2
cx2 = cx + c_w + cg
add_rect(slide, cx2, cy, c_w, c_h, BG_CARD)
add_text(slide, cx2 + Inches(0.2), cy + Inches(0.15), c_w - Inches(0.4), Inches(0.4),
         "② 정밀도 decay 강화", 16, ACCENT3, True)
add_bullet_text(slide, cx2 + Inches(0.2), cy + Inches(0.65), c_w - Inches(0.4), Inches(2.5), [
    ("k2: 0.3 → 0.5", WHITE, False),
    ("", WHITE, False),
    ("d_error 5m:  22 → 8.2  (3배↓)", WHITE, False),
    ("d_error 10m:  5 → 0.7  (7배↓)", WHITE, False),
    ("d_error 15m:  1 → 0.06 (17배↓)", WHITE, False),
    ("", WHITE, False),
    ("5m와 15m는", YELLOW, True),
    ("완전히 다른 결과라는 신호", YELLOW, True),
], 13, spacing=Pt(2))

# 변경 3
cx3 = cx2 + c_w + cg
add_rect(slide, cx3, cy, c_w, c_h, BG_CARD)
add_text(slide, cx3 + Inches(0.2), cy + Inches(0.15), c_w - Inches(0.4), Inches(0.4),
         "③ 예측 정확도 bonus 신규", 16, ACCENT3, True)
add_bullet_text(slide, cx3 + Inches(0.2), cy + Inches(0.65), c_w - Inches(0.4), Inches(2.5), [
    ("20 × exp(-0.1 × |gap|)", WHITE, False),
    ("", WHITE, False),
    ("gap = |d_impact − d_error|", GRAY, False),
    ("CCIP 예측과 실제 오차의 차이", GRAY, False),
    ("", WHITE, False),
    ("CCIP를 신뢰할 수 있는", WHITE, False),
    ("상황에서 투하하도록 유도", WHITE, False),
], 13, spacing=Pt(2))


# ════════════════════════════════════════════════════════════════
#  SLIDE 7: 보상 비교 테이블
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "투하 보상 — 이전 vs 현재 비교", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(6), Pt(3), ACCENT)

# 비교 테이블
data2 = [
    ["d_error", "이전 보상", "현재 보상", "변화"],
    ["0m (jackpot)", "350", "200", ""],
    ["1m", "224", "111", ""],
    ["5m", "172", "38", "4.5배 ↓"],
    ["10m", "155", "31", "5배 ↓"],
    ["15m (수렴점)", "151", "30", "5배 ↓"],
    ["50m (즉시 drop)", "150", "6", "25배 ↓"],
]
add_table(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5),
          7, 4, data2,
          col_widths=[Inches(3.2), Inches(2.2), Inches(2.2), Inches(2.4)])

# 핵심 메시지
add_rect(slide, Inches(2), Inches(5.3), Inches(9), Inches(1.5), BG_CARD, YELLOW)
add_bullet_text(slide, Inches(2.5), Inches(5.5), Inches(8), Inches(1.2), [
    ("이전:  최고(0m) vs 최저(50m)  =  350 vs 150  →  2.3배 차이 (둔감)", ACCENT2, True),
    ("현재:  최고(0m) vs 최저(50m)  =  200 vs 6      →  33배 차이 (민감)", ACCENT3, True),
], 17, spacing=Pt(10))


# ════════════════════════════════════════════════════════════════
#  SLIDE 8: 투하 메커니즘
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "투하 메커니즘 — Auto-only → Hybrid", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(6), Pt(3), ACCENT)

# 이전
add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5), BG_CARD, ACCENT2)
add_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.4),
         "이전: 자동 투하만 존재", 20, ACCENT2, True)
add_text(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.5),
         "d_impact ≤ 10m → 즉시 자동 투하", 16, WHITE, True, font_name="Consolas")
add_bullet_text(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(3.5), [
    ("에이전트가 투하 시점을 결정할 수 없음", WHITE, False),
    ("", WHITE, False),
    ("예측 낙하점이 10m 안에 처음 진입하는", WHITE, False),
    ("순간 강제로 drop이 발동", WHITE, False),
    ("", WHITE, False),
    ('"더 가까이 가서 정밀하게" 전략이', YELLOW, False),
    ("구조적으로 불가능", YELLOW, True),
    ("", WHITE, False),
    ("action[4] (투하 행동)가", GRAY, False),
    ("학습에 전혀 기여하지 못함", GRAY, False),
], 14, spacing=Pt(3))

# 현재
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), BG_CARD, ACCENT3)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.4),
         "현재: 하이브리드 투하", 20, ACCENT3, True)
add_bullet_text(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.8), [
    ("수동 = action[4] > 0.5  AND  step ≥ 10", WHITE, True),
    ("자동 = d_impact ≤ 2.0m  (안전망)", WHITE, True),
], 14, font_name="Consolas", spacing=Pt(4))
add_bullet_text(slide, Inches(7.1), Inches(3.2), Inches(5.2), Inches(3.5), [
    ("에이전트가 직접 투하 시점을 결정", ACCENT3, True),
    ("", WHITE, False),
    ("step ≥ 10: 이륙 직후 무의미 투하 방지", WHITE, False),
    ("자동 2m: 최소한의 안전망", WHITE, False),
    ("", WHITE, False),
    ("proximity bonus와 결합:", YELLOW, True),
    ("멀리서 즉시 투하 → 보상 ≈ 0", WHITE, False),
    ("가까이 가서 투하 → 보상 최대", WHITE, False),
    ("", WHITE, False),
    ("WandB에서 수동/자동 비율 추적", GRAY, False),
    ("→ 수동↑ = 에이전트가 학습 중인 증거", GRAY, False),
], 14, spacing=Pt(3))


# ════════════════════════════════════════════════════════════════
#  SLIDE 9: 종료 조건
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "에피소드 종료 조건 정비", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(5), Pt(3), ACCENT)

# 문제
add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.0), BG_CARD, ACCENT2)
add_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.4),
         "이전의 문제", 18, ACCENT2, True)
add_bullet_text(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(1.2), [
    ("지면에 닿아 바운스해도 종료되지 않음", WHITE, False),
    ("타겟 반대 방향으로 날아가도 500 step 대기", WHITE, False),
    ("→ 무의미한 에피소드로 학습 시간 낭비", YELLOW, True),
], 14, spacing=Pt(4))

# 테이블
data3 = [
    ["조건", "기준", "페널티", "비고"],
    ["ground contact", "고도 < 0.5m (무조건)", "-50", "신규 ★"],
    ["crash (저고도)", "고도 < 3m, step > 1", "-50", "유예 축소"],
    ["overspeed", "속도 > 20 m/s", "-30", ""],
    ["ang_vel", "각속도 > 2.0 rad/s", "-30", ""],
    ["inverted", "roll/pitch > 60°", "-30", ""],
    ["out_of_range", "타겟 거리 > 100m", "-30", "신규 ★"],
    ["timeout", "500 step (drop 미발동)", "-15", "페널티 축소"],
]
add_table(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(3.2),
          8, 4, data3,
          col_widths=[Inches(2.5), Inches(4.0), Inches(1.5), Inches(4.0)])

# 해결 카드
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.0), BG_CARD, ACCENT3)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.4),
         "추가된 조건", 18, ACCENT3, True)
add_bullet_text(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(1.2), [
    ("ground contact: 지면 바운스 제거", ACCENT3, False),
    ("out_of_range: 역주행 에피소드 조기 종료", ACCENT3, False),
    ("→ 학습 효율 향상", YELLOW, True),
], 14, spacing=Pt(4))


# ════════════════════════════════════════════════════════════════
#  SLIDE 10: 학습 설정 요약
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "학습 설정", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(3), Pt(3), ACCENT)

# SAC 테이블
data4 = [
    ["파라미터", "값", "비고"],
    ["learning_rate", "3e-4", "SAC 표준"],
    ["buffer_size", "500,000", "경험 보존 (100K→확대)"],
    ["batch_size", "256", "표준"],
    ["gamma", "0.995", "horizon ≈ 200 step"],
    ["gradient_steps", "1", "4에서 발산 → 1로 복구"],
    ["net_arch", "[256, 256]", "2-layer MLP"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(6.2), Inches(3.0),
          7, 3, data4,
          col_widths=[Inches(2.2), Inches(1.5), Inches(2.5)])

# 환경 테이블
data5 = [
    ["파라미터", "값", "비고"],
    ["action_vx_scale", "8 m/s", "15에서 축소"],
    ["action_rate_limit", "0.2", "급조작 방지"],
    ["max_steps", "500", "~25초"],
    ["min_altitude", "3.0m", "안전 고도"],
]
add_table(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(2.3),
          5, 3, data5,
          col_widths=[Inches(2.2), Inches(1.2), Inches(2.1)])

# 핵심 결정
add_rect(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.4), BG_CARD)
add_text(slide, Inches(0.9), Inches(4.9), Inches(5.2), Inches(0.4),
         "핵심 결정: gradient_steps = 1", 18, YELLOW, True)
add_bullet_text(slide, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.6), [
    ("4로 설정 시 critic loss 17,000 폭주", WHITE, False),
    ("큰 터미널 보상 + 작은 step 보상 혼재를", GRAY, False),
    ("4배 빈도로 학습 → Q값 발산", GRAY, False),
    ("1로 복귀 후 안정화 확인", ACCENT3, False),
], 14, spacing=Pt(3))

add_rect(slide, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.4), BG_CARD)
add_text(slide, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.4),
         "핵심 결정: gamma = 0.995", 18, YELLOW, True)
add_bullet_text(slide, Inches(7.1), Inches(5.4), Inches(5.2), Inches(1.6), [
    ("유효 horizon ≈ 200 step", WHITE, False),
    ("당장 다음 step이 아닌", GRAY, False),
    ("수십 초 후의 투하 성공까지 고려해", GRAY, False),
    ("비행 경로를 계획할 수 있게 된다", ACCENT3, False),
], 14, spacing=Pt(3))


# ════════════════════════════════════════════════════════════════
#  SLIDE 11: 현재 상태 & 향후
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "현재 상태와 향후 방향", 32, ACCENT, True)
add_rect(slide, Inches(0.6), Inches(0.95), Inches(5), Pt(3), ACCENT)

# 현재 상태
add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.4), BG_CARD, ACCENT)
add_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.4),
         "현재 진행 상황", 20, ACCENT, True)
add_bullet_text(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(1.5), [
    ("Round 1 본학습 진행 중", WHITE, True),
    ("WandB run: dv428atu", GRAY, False),
    ("150k step, fresh start", GRAY, False),
    ("이전 경험 버퍼 초기화 후 시작", GRAY, False),
], 15, spacing=Pt(3))

# 확인 지표
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.4), BG_CARD)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.4),
         "확인할 지표", 20, YELLOW, True)
add_bullet_text(slide, Inches(7.1), Inches(2.0), Inches(5.2), Inches(1.5), [
    ("drop_error:  13m → 목표 <5m", WHITE, False),
    ("total_success:  0회 → 1회 이상", WHITE, False),
    ("수동 투하 비율:  에이전트 학습 증거", WHITE, False),
    ("critic_loss:  <300 안정 유지", WHITE, False),
], 15, spacing=Pt(3))

# 결과별 분기
add_rect(slide, Inches(0.6), Inches(4.2), Inches(3.7), Inches(2.8), BG_CARD, ACCENT3)
add_text(slide, Inches(0.9), Inches(4.3), Inches(3.1), Inches(0.4),
         "성공 시", 18, ACCENT3, True)
add_bullet_text(slide, Inches(0.9), Inches(4.8), Inches(3.1), Inches(2.0), [
    ("Phase 2 진입", ACCENT3, True),
    ("auto_drop 축소/제거", WHITE, False),
    ("curriculum 고도화", WHITE, False),
], 14, spacing=Pt(3))

add_rect(slide, Inches(4.7), Inches(4.2), Inches(3.7), Inches(2.8), BG_CARD, YELLOW)
add_text(slide, Inches(5.0), Inches(4.3), Inches(3.1), Inches(0.4),
         "부분 개선 시", 18, YELLOW, True)
add_bullet_text(slide, Inches(5.0), Inches(4.8), Inches(3.1), Inches(2.0), [
    ("보류 안건 추가 적용", YELLOW, True),
    ("이중 보상 (#010)", WHITE, False),
    ("prediction bonus 조정", WHITE, False),
], 14, spacing=Pt(3))

add_rect(slide, Inches(8.8), Inches(4.2), Inches(3.7), Inches(2.8), BG_CARD, ACCENT2)
add_text(slide, Inches(9.1), Inches(4.3), Inches(3.1), Inches(0.4),
         "개선 없을 시", 18, ACCENT2, True)
add_bullet_text(slide, Inches(9.1), Inches(4.8), Inches(3.1), Inches(2.0), [
    ("근본 재검토", ACCENT2, True),
    ("관측 공간 변경", WHITE, False),
    ("알고리즘 교체 검토", WHITE, False),
], 14, spacing=Pt(3))


# ════════════════════════════════════════════════════════════════
#  SLIDE 12: 마무리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "정리", 36, ACCENT, True, PP_ALIGN.CENTER)

add_rect(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(3.5), BG_CARD)
add_bullet_text(slide, Inches(2.0), Inches(3.0), Inches(9), Inches(3.0), [
    ("7회 학습을 통해 인프라, 알고리즘, 보상 구조의 문제를 단계적으로 진단", WHITE, False),
    ("", WHITE, False),
    ("Round 1 핵심 변경:", YELLOW, True),
    ("  보상 재설계 — 정밀도에 민감하게, 스케일은 적정 범위로", WHITE, False),
    ("  하이브리드 투하 — 에이전트가 직접 타이밍 결정", WHITE, False),
    ("  종료 조건 정비 — 무의미한 에피소드 조기 종료", WHITE, False),
    ("", WHITE, False),
    ("현재 Round 1 본학습 진행 중 — 결과에 따라 다음 Phase 결정", ACCENT, True),
], 18, spacing=Pt(8))

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Phase 1 → Phase 2 (정밀화) → Phase 3 (일반화)  |  최종: 0.1m 이내 정밀투하", 16, GRAY, False, PP_ALIGN.CENTER)


# ── 저장 ──
out_path = "/home/juns/Drone-Bombard-Simulation/local/presentation/round1_presentation.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}")
